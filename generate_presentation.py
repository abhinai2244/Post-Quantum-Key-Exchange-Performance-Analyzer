from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Initialize presentation
prs = Presentation()

# Apply a simple template (0 is title, 1 is title and content)
title_slide_layout = prs.slide_layouts[0]
content_slide_layout = prs.slide_layouts[1]

# Helper function to style titles
def style_title(title_shape):
    if title_shape.text_frame.paragraphs:
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204) # Deep blue
        title_shape.text_frame.paragraphs[0].font.bold = True

# ==========================================
# Slide 1: Title & Problem Statement
# ==========================================
slide1 = prs.slides.add_slide(title_slide_layout)
title1 = slide1.shapes.title
subtitle1 = slide1.placeholders[1]

title1.text = "Post-Quantum Key Exchange Performance Analyzer"
subtitle1.text = "Visualizing the Fall of Classical Cryptography\nand the Rise of PQC\n\nHackathon Submission"

# ==========================================
# Slide 2: Objectives
# ==========================================
slide2 = prs.slides.add_slide(content_slide_layout)
title2 = slide2.shapes.title
body2 = slide2.placeholders[1]

title2.text = "Objectives"
style_title(title2)

tf2 = body2.text_frame
tf2.text = "Demonstrate the mathematical vulnerability of classical cryptography to quantum computers."

p = tf2.add_paragraph()
p.text = "Simulate small-scale quantum attacks:"
p.level = 1

p = tf2.add_paragraph()
p.text = "Shor's Algorithm against RSA (Prime Factorization)"
p.level = 2

p = tf2.add_paragraph()
p.text = "Grover's Algorithm against AES (Search Space Amplification)"
p.level = 2

p = tf2.add_paragraph()
p.text = "Analyze hardware limitations by injecting real-world NISQ-era noise."
p.level = 1

p = tf2.add_paragraph()
p.text = "Educate and visualize why Post-Quantum schemes (Lattice-based Kyber & Hash-based SPHINCS+) resist these quantum threats."
p.level = 1


# ==========================================
# Slide 3: Implementation Details
# ==========================================
slide3 = prs.slides.add_slide(content_slide_layout)
title3 = slide3.shapes.title
body3 = slide3.placeholders[1]

title3.text = "Implementation Details"
style_title(title3)

tf3 = body3.text_frame
tf3.text = "Built strictly as a robust, modular Python architecture."

p = tf3.add_paragraph()
p.text = "Backend Quantum Engine:"
p.level = 1

p = tf3.add_paragraph()
p.text = "IBM Qiskit Aer Simulator for parameterized Shor/Grover circuits."
p.level = 2

p = tf3.add_paragraph()
p.text = "NoiseModule mapping Depolarizing, Phase, and Bit-flip errors."
p.level = 2

p = tf3.add_paragraph()
p.text = "Interactive Dashboard (Streamlit & Plotly):"
p.level = 1

p = tf3.add_paragraph()
p.text = "Dark Glassmorphism UI ensuring maximum visual impact."
p.level = 2

p = tf3.add_paragraph()
p.text = "7 dynamic tabs featuring live simulation races and 3D heatmaps."
p.level = 2

p = tf3.add_paragraph()
p.text = "Quality Assurance: 100% Pytest coverage mathematically validating cryptographic integrity."
p.level = 1


# ==========================================
# Slide 4: Results (Graphs/Charts Context)
# ==========================================
slide4 = prs.slides.add_slide(content_slide_layout)
title4 = slide4.shapes.title
body4 = slide4.placeholders[1]

title4.text = "Results & Visualizations"
style_title(title4)

tf4 = body4.text_frame
tf4.text = "The analyzer produces real-time, interactive insights:"

p = tf4.add_paragraph()
p.text = "Live Attack Race: Shows Quantum O(log N)³ vs Classical Sub-exponential speeds."
p.level = 1

p = tf4.add_paragraph()
p.text = "NISQ Noise 3D Surface Plot:"
p.level = 1

p = tf4.add_paragraph()
p.text = "Proves that >3% hardware noise destroys quantum advantage, effectively returning Grover's algorithm to blind random guessing."
p.level = 2

p = tf4.add_paragraph()
p.text = "Algorithm Security Radar:"
p.level = 1

p = tf4.add_paragraph()
p.text = "Directly compares RSA, AES, Kyber, and SPHINCS+ across key-size efficiency, standardization maturity, and quantum resistance."
p.level = 2


# ==========================================
# Slide 5: Analysis, Conclusions, and Future Work
# ==========================================
slide5 = prs.slides.add_slide(content_slide_layout)
title5 = slide5.shapes.title
body5 = slide5.placeholders[1]

title5.text = "Analysis & Conclusions"
style_title(title5)

tf5 = body5.text_frame
tf5.text = "Analysis Highlights:"

p = tf5.add_paragraph()
p.text = "Classical cryptography is polynomial-time fragile to quantum algorithms."
p.level = 1

p = tf5.add_paragraph()
p.text = "Post-Quantum schemes relying on Lattice structures (LWE) and extended Hash boundaries survive natively."
p.level = 1

p = tf5.add_paragraph()
p.text = "Future Work:"
p.level = 1

p = tf5.add_paragraph()
p.text = "Extend Qiskit backend to execute on real IBM Quantum Hardware via cloud API."
p.level = 2

p = tf5.add_paragraph()
p.text = "Integrate side-channel attack vulnerabilities to the Post-Quantum algorithms."
p.level = 2


# Save the presentation
output_path = "Post_Quantum_Analyzer_Presentation.pptx"
prs.save(output_path)
print(f"Presentation successfully generated at: {output_path}")
